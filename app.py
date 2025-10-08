import io
import zipfile
from pathlib import Path
from itertools import count
import time

import numpy as np
import fitz  # PyMuPDF
import cv2
from PIL import Image
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt

# ========================= Sabit Ayarlar =========================
DPI = 240
ROW_RATIO_THRESH = 0.45
CNT_MIN_W_RATIO = 0.65
MIN_BAND_H = 32
MAX_BAND_H_FRAC = 1 / 10
PAD_TOP = 6
PAD_BOTTOM = 12


# ========================= Yardımcılar =========================

def st_image_compat(col, img, caption=None):
    """Streamlit sürüm farkı için uyumlu görüntüleme."""
    try:
        col.image(img, caption=caption, width='stretch')
    except TypeError:
        try:
            col.image(img, caption=caption, use_container_width=True)
        except TypeError:
            col.image(img, caption=caption, use_column_width=True)


def pdf_to_images(file_bytes: bytes, dpi: int = DPI):
    """PDF dosyasını (bytes) sayfa görsellerine çevirir (BGR numpy)."""
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    imgs = []
    zoom = dpi / 72.0
    mat = fitz.Matrix(zoom, zoom)
    for page in doc:
        pix = page.get_pixmap(matrix=mat, alpha=False)
        img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, 3)
        img = cv2.cvtColor(img, cv2.COLOR_RGB2BGR)
        imgs.append(img)
    doc.close()
    return imgs


def red_mask(bgr_img, lower1, upper1, lower2, upper2):
    hsv = cv2.cvtColor(bgr_img, cv2.COLOR_BGR2HSV)
    m1 = cv2.inRange(hsv, np.array(lower1), np.array(upper1))
    m2 = cv2.inRange(hsv, np.array(lower2), np.array(upper2))
    mask = cv2.bitwise_or(m1, m2)
    kernel = np.ones((5, 5), np.uint8)
    mask = cv2.morphologyEx(mask, cv2.MORPH_CLOSE, kernel, iterations=2)
    return mask


def headers_by_row(mask, ratio_thresh, min_h, max_h, gap_tol):
    """Satır projeksiyonu ile geniş kırmızı şeritleri bul."""
    H, W = mask.shape[:2]
    red_ratio = mask.sum(axis=1) / (255.0 * W)
    win = max(7, H // 200)
    kernel1d = np.ones(win, dtype=np.float32) / win
    smooth = np.convolve(red_ratio, kernel1d, mode='same')

    bands, in_run, run_start, last_y = [], False, 0, -10 ** 9
    for y, val in enumerate(smooth):
        if val >= ratio_thresh:
            if not in_run:
                if bands and (y - bands[-1][1] <= gap_tol):
                    in_run, run_start = True, bands[-1][0]
                    bands.pop()
                else:
                    in_run, run_start = True, y
            last_y = y
        else:
            if in_run:
                in_run = False
                if last_y - run_start + 1 >= min_h:
                    bands.append((run_start, last_y))
    if in_run and (last_y - run_start + 1 >= min_h):
        bands.append((run_start, last_y))

    cleaned = []
    for s, e in bands:
        h = e - s + 1
        if h > max_h:
            mid = (s + e) // 2
            half = max_h // 2
            cleaned.append((max(0, mid - half), min(H - 1, mid + half)))
        else:
            cleaned.append((s, e))
    return cleaned


def headers_by_contours(mask, w_ratio, min_h, max_h):
    """Kontur tabanlı: tam genişliğe yakın yatay kırmızı şeritler."""
    H, W = mask.shape[:2]
    cnts, _ = cv2.findContours(mask, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    out = []
    for c in cnts:
        x, y, w, h = cv2.boundingRect(c)
        if w >= w_ratio * W and (min_h <= h <= max_h):
            out.append((y, y + h - 1))
    return out


def merge_bands(bands, gap_tol):
    if not bands:
        return []
    bands = sorted(bands, key=lambda t: t[0])
    merged = []
    cs, ce = bands[0]
    for s, e in bands[1:]:
        if s <= ce + gap_tol:
            ce = max(ce, e)
        else:
            merged.append((cs, ce))
            cs, ce = s, e
    merged.append((cs, ce))
    return merged


def find_header_bands(img,
                      row_ratio_thresh=ROW_RATIO_THRESH,
                      min_band_h=MIN_BAND_H,
                      max_band_h_frac=MAX_BAND_H_FRAC,
                      gap_frac=1 / 300,
                      cnt_min_w_ratio=CNT_MIN_W_RATIO,
                      cnt_min_h=28,
                      cnt_max_h=260,
                      hsv1_low=(0, 80, 70),
                      hsv1_up=(12, 255, 255),
                      hsv2_low=(170, 80, 70),
                      hsv2_up=(180, 255, 255)):
    """Şeritleri iki yöntemle bul, birleştir."""
    H, W = img.shape[:2]
    gap_tol = max(6, int(H * gap_frac))
    max_band_h = max(int(H * max_band_h_frac), min_band_h + 10)

    mask = red_mask(img, hsv1_low, hsv1_up, hsv2_low, hsv2_up)
    b1 = headers_by_row(mask, row_ratio_thresh, min_band_h, max_band_h, gap_tol)
    b2 = headers_by_contours(mask, cnt_min_w_ratio, cnt_min_h, cnt_max_h)
    bands = merge_bands(b1 + b2, gap_tol)
    return bands, mask


def slice_by_headers(img, header_bands, pad_top=PAD_TOP, pad_bottom=PAD_BOTTOM):
    H, W = img.shape[:2]
    if not header_bands:
        return [(0, 0, W, H)]
    header_bands = sorted(header_bands, key=lambda t: t[0])
    crops = []
    for i in range(len(header_bands)):
        y0 = max(0, header_bands[i][0] - pad_top)
        if i < len(header_bands) - 1:
            y1 = header_bands[i + 1][0] - 1
        else:
            y1 = H
        y1 = min(H, y1 + pad_bottom)
        crops.append((0, y0, W, y1))
    return crops


def bgr_to_pil(bgr):
    rgb = cv2.cvtColor(bgr, cv2.COLOR_BGR2RGB)
    return Image.fromarray(rgb)


def create_powerpoint(images_bytes_list):
    """Görselleri PowerPoint sunumuna dönüştür (16:9 optimize edilmiş)."""
    prs = Presentation()

    # 16:9 oran ayarla (standart slide boyutu)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    for img_bytes in images_bytes_list:
        # Boş slide ekle
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)

        # Görseli PIL olarak aç
        img = Image.open(io.BytesIO(img_bytes))

        # Görsel boyutlarını al
        img_width, img_height = img.size
        img_ratio = img_width / img_height

        # Slide boyutları
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        slide_ratio = slide_width / slide_height

        # Görseli slide'a sığdır (aspect ratio koruyarak)
        if img_ratio > slide_ratio:
            # Görsel daha geniş - genişliğe göre ayarla
            pic_width = slide_width
            pic_height = int(slide_width / img_ratio)
        else:
            # Görsel daha uzun - yüksekliğe göre ayarla
            pic_height = slide_height
            pic_width = int(slide_height * img_ratio)

        # Görseli ortala
        left = (slide_width - pic_width) // 2
        top = (slide_height - pic_height) // 2

        # Geçici bytes buffer'a kaydet
        img_buffer = io.BytesIO(img_bytes)

        # Slide'a ekle
        slide.shapes.add_picture(img_buffer, left, top, width=pic_width, height=pic_height)

    # PowerPoint'i bytes olarak döndür
    pptx_buffer = io.BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer


# ========================= Arayüz =========================

st.set_page_config(
    page_title="PDF Slide Kesici",
    page_icon="✂️",
    layout="wide",
    initial_sidebar_state="collapsed"
)


# Kalıcı durum
def _ensure_state():
    if "crops_pngs" not in st.session_state:
        st.session_state.crops_pngs = []
    if "last_count" not in st.session_state:
        st.session_state.last_count = 0
    if "processed" not in st.session_state:
        st.session_state.processed = False
    if "pptx_created" not in st.session_state:
        st.session_state.pptx_created = False
    if "pptx_buffer" not in st.session_state:
        st.session_state.pptx_buffer = None
    if "show_preview" not in st.session_state:
        st.session_state.show_preview = False
    if "preview_timestamp" not in st.session_state:
        st.session_state.preview_timestamp = None
    if "downloaded" not in st.session_state:
        st.session_state.downloaded = False


_ensure_state()


# Otomatik temizlik kontrolü (5 dakika)
def check_auto_cleanup():
    if st.session_state.show_preview and st.session_state.preview_timestamp:
        elapsed = time.time() - st.session_state.preview_timestamp
        if elapsed > 300:  # 5 dakika = 300 saniye
            st.session_state.show_preview = False
            st.session_state.crops_pngs = []
            st.session_state.processed = False
            st.session_state.pptx_created = False
            st.session_state.pptx_buffer = None
            st.session_state.downloaded = False
            st.rerun()


check_auto_cleanup()

# ---- Modern CSS Stilleri ----
st.markdown(
    """
    <style>
    /* Ana başlık stili */
    h1 {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        font-size: 2.5rem !important;
        font-weight: 800 !important;
        margin-bottom: 0.5rem !important;
    }

    /* Kart stilleri */
    .upload-card {
        background: linear-gradient(135deg, #667eea15 0%, #764ba215 100%);
        border: 2px dashed #667eea;
        border-radius: 20px;
        padding: 3rem 2rem;
        text-align: center;
        transition: all 0.3s ease;
        margin: 2rem 0;
    }

    .upload-card:hover {
        border-color: #764ba2;
        box-shadow: 0 8px 25px rgba(102, 126, 234, 0.2);
        transform: translateY(-2px);
    }

    /* Buton stilleri */
    .stDownloadButton button {
        background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%) !important;
        color: white !important;
        border: none !important;
        border-radius: 12px !important;
        font-weight: 700 !important;
        padding: 0.7rem 2rem !important;
        box-shadow: 0 6px 20px rgba(245, 87, 108, 0.35) !important;
        transition: all 0.3s ease !important;
        width: 100% !important;
    }

    .stDownloadButton button:hover {
        transform: translateY(-2px) !important;
        box-shadow: 0 8px 25px rgba(245, 87, 108, 0.45) !important;
    }

    /* İkincil buton (Slide'lara Aktar) */
    .stButton button[kind="secondary"] {
        background: linear-gradient(135deg, #4facfe 0%, #00f2fe 100%) !important;
        color: white !important;
        border: none !important;
        border-radius: 12px !important;
        font-weight: 700 !important;
        padding: 0.7rem 2rem !important;
        box-shadow: 0 6px 20px rgba(79, 172, 254, 0.35) !important;
        transition: all 0.3s ease !important;
        width: 100% !important;
    }

    .stButton button[kind="secondary"]:hover {
        transform: translateY(-2px) !important;
        box-shadow: 0 8px 25px rgba(79, 172, 254, 0.45) !important;
    }
    .stButton button[kind="primary"] {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%) !important;
        border: none !important;
        border-radius: 12px !important;
        font-weight: 700 !important;
        padding: 0.7rem 2rem !important;
        box-shadow: 0 6px 20px rgba(102, 126, 234, 0.35) !important;
        transition: all 0.3s ease !important;
        width: 100% !important;
    }

    .stButton button[kind="primary"]:hover {
        transform: translateY(-2px) !important;
        box-shadow: 0 8px 25px rgba(102, 126, 234, 0.45) !important;
    }

    /* İstatistik kartları */
    .stat-card {
        background: white;
        border-radius: 16px;
        padding: 1.5rem;
        box-shadow: 0 4px 15px rgba(0, 0, 0, 0.08);
        text-align: center;
        border-left: 4px solid;
        transition: all 0.3s ease;
    }

    .stat-card:hover {
        transform: translateY(-3px);
        box-shadow: 0 8px 25px rgba(0, 0, 0, 0.12);
    }

    .stat-number {
        font-size: 2.5rem;
        font-weight: 800;
        margin: 0.5rem 0;
    }

    .stat-label {
        color: #666;
        font-size: 0.95rem;
        font-weight: 600;
        text-transform: uppercase;
        letter-spacing: 0.5px;
    }

    /* Önizleme kartları */
    [data-testid="stImage"] {
        border-radius: 12px;
        overflow: hidden;
        box-shadow: 0 4px 15px rgba(0, 0, 0, 0.1);
        transition: all 0.3s ease;
    }

    [data-testid="stImage"]:hover {
        transform: scale(1.02);
        box-shadow: 0 8px 25px rgba(0, 0, 0, 0.15);
    }

    /* Info kutusu */
    .stAlert {
        border-radius: 12px;
        border-left: 4px solid #667eea;
    }

    /* Başarı mesajı */
    .success-message {
        background: linear-gradient(135deg, #84fab015 0%, #8fd3f415 100%);
        border: 2px solid #84fab0;
        border-radius: 16px;
        padding: 1.5rem;
        margin: 2rem 0;
        text-align: center;
    }

    /* File uploader stili */
    [data-testid="stFileUploader"] {
        border-radius: 16px;
    }

    [data-testid="stFileUploader"] > div {
        border-radius: 16px;
    }
    </style>
    """,
    unsafe_allow_html=True,
)

# ---- Başlık ----
st.markdown("<h1>✂️ PDF Slide Kesici</h1>", unsafe_allow_html=True)
st.markdown("### 🎯 Kırmızı başlık şeritlerine göre otomatik kesim")
st.markdown("---")

# ---- Dosya Yükleme Bölümü ----
col1, col2, col3 = st.columns([1, 2, 1])

with col2:
    st.markdown("### 📁 PDF Dosyanızı Yükleyin")
    uploaded = st.file_uploader(
        "PDF dosyasını sürükleyin veya seçin",
        type=["pdf"],
        accept_multiple_files=False,
        label_visibility="collapsed"
    )

    if uploaded:
        file_size = len(uploaded.getvalue()) / (1024 * 1024)
        st.success(f"✅ **{uploaded.name}** yüklendi ({file_size:.1f} MB)")

        st.markdown("<br>", unsafe_allow_html=True)

        # İşle butonu - tam genişlik
        if st.button("🚀 İşlemeyi Başlat", type="primary", key="process_btn"):
            pdf_bytes = uploaded.getvalue()

            # Progress bar ile işleme
            progress_bar = st.progress(0)
            status_text = st.empty()

            status_text.text("📄 PDF rasterize ediliyor...")
            progress_bar.progress(20)
            pages = pdf_to_images(pdf_bytes, dpi=DPI)

            crops_pngs = []
            total_bands = 0

            status_text.text("🔍 Başlık tespiti ve kırpma yapılıyor...")
            progress_bar.progress(40)

            for idx, img in enumerate(pages):
                bands, _ = find_header_bands(img)
                total_bands += len(bands)
                boxes = slice_by_headers(img, bands)
                for box in boxes:
                    x0, y0, x1, y1 = box
                    crop = img[y0:y1, x0:x1]
                    pil_im = bgr_to_pil(crop)
                    b = io.BytesIO()
                    pil_im.save(b, format="PNG")
                    crops_pngs.append(b.getvalue())

                progress_bar.progress(40 + int(50 * (idx + 1) / len(pages)))

            status_text.text("✨ Tamamlanıyor...")
            progress_bar.progress(100)

            # Duruma kaydet
            st.session_state.crops_pngs = crops_pngs
            st.session_state.last_count = len(crops_pngs)
            st.session_state.processed = True
            st.session_state.pages_count = len(pages)
            st.session_state.bands_count = total_bands

            status_text.empty()
            progress_bar.empty()

            st.rerun()

st.markdown("---")

# ---- İstatistikler ve İndirme ----
if st.session_state.processed and st.session_state.crops_pngs:
    # İstatistik kartları
    st.markdown("### 📊 İşlem Sonuçları")
    col1, col2, col3, col4 = st.columns(4)

    with col1:
        st.markdown(
            f"""
            <div class="stat-card" style="border-left-color: #667eea;">
                <div class="stat-label">📄 Sayfa</div>
                <div class="stat-number" style="color: #667eea;">{st.session_state.pages_count}</div>
            </div>
            """,
            unsafe_allow_html=True
        )

    with col2:
        st.markdown(
            f"""
            <div class="stat-card" style="border-left-color: #f5576c;">
                <div class="stat-label">🎯 Başlık</div>
                <div class="stat-number" style="color: #f5576c;">{st.session_state.bands_count}</div>
            </div>
            """,
            unsafe_allow_html=True
        )

    with col3:
        st.markdown(
            f"""
            <div class="stat-card" style="border-left-color: #4facfe;">
                <div class="stat-label">✂️ Kesim</div>
                <div class="stat-number" style="color: #4facfe;">{st.session_state.last_count}</div>
            </div>
            """,
            unsafe_allow_html=True
        )

    with col4:
        # ZIP indirme butonu
        zip_buf = io.BytesIO()
        with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
            for i, data in enumerate(st.session_state.crops_pngs, start=1):
                zf.writestr(f"slide_{i:04d}.png", data)
        zip_buf.seek(0)

        if st.download_button(
                label=f"📦 ZIP İndir\n({st.session_state.last_count} görsel)",
                data=zip_buf,
                file_name="slides.zip",
                mime="application/zip",
                key="download_zip"
        ):
            # İndirme yapıldı, önizlemeyi temizle
            st.session_state.show_preview = False
            st.session_state.crops_pngs = []
            st.session_state.processed = False
            st.session_state.downloaded = True
            time.sleep(0.5)  # Kısa gecikme
            st.rerun()

    st.markdown("<br>", unsafe_allow_html=True)

    # PowerPoint oluşturma bölümü
    col_ppt1, col_ppt2, col_ppt3 = st.columns([1, 2, 1])

    with col_ppt2:
        if not st.session_state.pptx_created:
            if st.button("🎯 Slide'lara Aktar (PowerPoint Oluştur)", type="secondary", key="create_ppt"):
                with st.spinner("📊 PowerPoint sunumu oluşturuluyor..."):
                    pptx_buffer = create_powerpoint(st.session_state.crops_pngs)
                    st.session_state.pptx_buffer = pptx_buffer
                    st.session_state.pptx_created = True
                    st.success(f"✅ PowerPoint sunumu hazır! {st.session_state.last_count} slide oluşturuldu.")
                    st.rerun()
        else:
            st.success(f"✅ PowerPoint hazır! {st.session_state.last_count} slide içeriyor.")

            # PPTX indirme butonu
            if st.download_button(
                    label="📥 PowerPoint'i İndir (.pptx)",
                    data=st.session_state.pptx_buffer,
                    file_name="sunum.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    type="primary",
                    key="download_ppt"
            ):
                # İndirme yapıldı, önizlemeyi temizle
                st.session_state.show_preview = False
                st.session_state.crops_pngs = []
                st.session_state.processed = False
                st.session_state.pptx_created = False
                st.session_state.pptx_buffer = None
                st.session_state.downloaded = True
                time.sleep(0.5)  # Kısa gecikme
                st.rerun()

    st.markdown("<br>", unsafe_allow_html=True)
    st.markdown("---")

    # Önizleme bölümü - sadece show_preview True ise göster
    if st.session_state.show_preview:
        # Kalan süreyi hesapla
        if st.session_state.preview_timestamp:
            elapsed = time.time() - st.session_state.preview_timestamp
            remaining = max(0, 300 - int(elapsed))  # 5 dakika = 300 saniye
            minutes = remaining // 60
            seconds = remaining % 60

        st.markdown("### 🖼️ Kesilen Görseller")

        col_info1, col_info2, col_info3 = st.columns([1, 1, 1])
        with col_info1:
            st.info(f"⏱️ Otomatik temizlik: {minutes}:{seconds:02d}")
        with col_info2:
            st.caption("ZIP veya PPTX indirince otomatik silinir")
        with col_info3:
            if st.button("🗑️ Şimdi Temizle", key="clear_preview"):
                st.session_state.show_preview = False
                st.session_state.crops_pngs = []
                st.session_state.processed = False
                st.rerun()

        st.caption("Görsellerin üzerine gelerek büyütebilirsiniz")

        # Izgara görünümü - 3 sütun
        for i in range(0, len(st.session_state.crops_pngs), 3):
            cols = st.columns(3)
            for j in range(3):
                idx = i + j
                if idx < len(st.session_state.crops_pngs):
                    data = st.session_state.crops_pngs[idx]
                    im = Image.open(io.BytesIO(data)).convert("RGB")
                    with cols[j]:
                        st_image_compat(cols[j], im, caption=f"Slide {idx + 1:04d}")

        # Her 10 saniyede bir yenile (countdown için)
        time.sleep(0.1)
        st.rerun()

else:
    # Boş durum mesajı
    if st.session_state.downloaded:
        st.success("✅ Dosyanız indirildi! Önizleme temizlendi.")
        st.info("👆 Yeni bir PDF yükleyerek tekrar işlem yapabilirsiniz", icon="ℹ️")
    else:
        st.info("👆 PDF dosyanızı yükleyin ve işleme başlatın", icon="ℹ️")

    # Nasıl çalışır bölümü
    with st.expander("❓ Nasıl Çalışır?"):
        st.markdown("""
        **Bu araç neler yapar?**

        1. 📤 PDF dosyanızı yüklersiniz
        2. 🔍 Sistem kırmızı başlık şeritlerini otomatik tespit eder
        3. ✂️ Her başlık arasını ayrı görsel olarak keser
        4. 📦 Tüm görselleri ZIP dosyası olarak indirebilirsiniz

        **İdeal kullanım alanları:**
        - Sunum slaytlarını ayırmak
        - Ders notlarını bölümlere ayırmak
        - Rapor sayfalarını bireysel görsellere dönüştürmek

        **Teknik detaylar:**
        - DPI: 240 (yüksek kalite)
        - Format: PNG (görsel) / PPTX (sunum)
        - Otomatik padding ve hizalama
        - 16:9 oranında PowerPoint slayları
        """)

    with st.expander("⚙️ Özellikler"):
        st.markdown("""
        - ✨ Otomatik kırmızı renk algılama
        - 🎯 Akıllı şerit birleştirme
        - 📏 Otomatik padding ayarı
        - 🖼️ Yüksek kaliteli çıktı
        - 📦 Toplu indirme desteği
        - 🎯 PowerPoint sunumu oluşturma
        - 📐 Otomatik slide optimizasyonu
        - 👁️ Canlı önizleme
        """)

# Footer
st.markdown("<br><br>", unsafe_allow_html=True)
st.markdown("---")
st.markdown(
    "<p style='text-align: center; color: #999; font-size: 0.9rem;'>Made with ❤️ using Streamlit</p>",
    unsafe_allow_html=True
)